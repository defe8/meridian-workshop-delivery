"""Generate capabilities deck for Meridian Components RFP MC-2026-0417."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SLATE_DARK = RGBColor(0x0F, 0x17, 0x2A)
SLATE_MID = RGBColor(0x47, 0x55, 0x69)
SLATE_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT = RGBColor(0xD9, 0x77, 0x06)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=SLATE_DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_bullets(slide, left, top, width, height, bullets, *, size=18, color=SLATE_DARK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def add_accent_bar(slide, left, top, width=0.08, height=0.6, color=ACCENT):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color


def title_slide(prs, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, SLATE_DARK)
    add_accent_bar(slide, 0.7, 2.6, width=0.8, height=0.08)
    add_text(slide, 0.7, 2.8, 12, 1.5, title, size=54, bold=True, color=WHITE)
    add_text(slide, 0.7, 4.2, 12, 1.0, subtitle, size=24, color=SLATE_LIGHT)
    add_text(slide, 0.7, 6.7, 12, 0.4, footer, size=12, color=SLATE_LIGHT)


def content_slide(prs, eyebrow, title, bullets, *, body=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_accent_bar(slide, 0.7, 0.7, height=0.35)
    add_text(slide, 0.95, 0.65, 10, 0.4, eyebrow.upper(), size=12, bold=True, color=ACCENT)
    add_text(slide, 0.7, 1.15, 12, 1.0, title, size=36, bold=True, color=SLATE_DARK)
    if body:
        add_text(slide, 0.7, 2.2, 12, 1.0, body, size=18, color=SLATE_MID)
        add_bullets(slide, 0.7, 3.0, 12, 4, bullets, size=18, color=SLATE_DARK)
    else:
        add_bullets(slide, 0.7, 2.4, 12, 4.5, bullets, size=20, color=SLATE_DARK)


def two_col_slide(prs, eyebrow, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_accent_bar(slide, 0.7, 0.7, height=0.35)
    add_text(slide, 0.95, 0.65, 10, 0.4, eyebrow.upper(), size=12, bold=True, color=ACCENT)
    add_text(slide, 0.7, 1.15, 12, 1.0, title, size=36, bold=True, color=SLATE_DARK)
    add_text(slide, 0.7, 2.4, 6, 0.5, left_title, size=18, bold=True, color=SLATE_DARK)
    add_bullets(slide, 0.7, 3.0, 5.8, 4, left_bullets, size=15, color=SLATE_MID)
    add_text(slide, 7.0, 2.4, 6, 0.5, right_title, size=18, bold=True, color=SLATE_DARK)
    add_bullets(slide, 7.0, 3.0, 5.8, 4, right_bullets, size=15, color=SLATE_MID)


def closing_slide(prs, title, subtitle, contact):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, SLATE_DARK)
    add_accent_bar(slide, 0.7, 2.6, width=0.8, height=0.08)
    add_text(slide, 0.7, 2.8, 12, 1.5, title, size=48, bold=True, color=WHITE)
    add_text(slide, 0.7, 4.0, 12, 1.0, subtitle, size=22, color=SLATE_LIGHT)
    add_text(slide, 0.7, 6.7, 12, 0.4, contact, size=12, color=SLATE_LIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(
        prs,
        "Modernizing Meridian's\nInventory Dashboard",
        "Capabilities Presentation  |  RFP MC-2026-0417",
        "Prepared for Meridian Components, Inc.  •  April 27, 2026",
    )

    content_slide(
        prs,
        "Act 1 — Understanding",
        "What we heard from Meridian",
        [
            "Operations runs the business through this dashboard — every day, three warehouses",
            "Reports module has known defects the previous vendor never closed out",
            "Restocking workflow VP Operations has been asking for doesn't exist yet",
            "No automated tests — IT has frozen further change until that's resolved",
            "The system works, but it can't safely evolve. That's the gap we close.",
        ],
    )

    content_slide(
        prs,
        "Scope",
        "Required deliverables",
        [
            "R1  Reports module remediation — audit + fix the 8+ logged defects",
            "R2  Restocking recommendations — budget-aware PO engine, operator-in-control",
            "R3  Automated browser tests — coverage on the flows that matter",
            "R4  Architecture documentation — IT-ready handoff, week one",
        ],
    )

    content_slide(
        prs,
        "Scope",
        "Desired — where they fit naturally",
        [
            "D1  UI refresh — modernize layer above your existing slate/gray tokens",
            "D2  Internationalization — Japanese for Tokyo, scaffold for future locales",
            "D3  Dark mode — operator-selectable, useful on warehouse floor stations",
            "We bring these in during Phase 3 only if Phase 1–2 land clean.",
        ],
    )

    content_slide(
        prs,
        "How we work",
        "Our approach",
        [
            "Tests aren't a deliverable at the end — they're how we work each week",
            "Architecture review week one, before we touch a line of code",
            "Restocking is designed around the operator, algorithm proposes, human decides",
            "Weekly demos with R. Tanaka. Defect log shared from day one.",
            "Predictable cadence. No surprises in week eight.",
        ],
    )

    two_col_slide(
        prs,
        "Timeline",
        "10 weeks, three phases",
        "Phase 1 — Foundation (wks 1–3)",
        [
            "Architecture review (R4)",
            "Reports audit & remediation (R1)",
            "First E2E test suite (R3)",
            "Milestone: green build, IT sign-off",
        ],
        "Phase 2 — Build (wks 4–7)",
        [
            "Restocking recommendation engine (R2)",
            "Budget ceiling + demand forecast",
            "Operator override workflow",
            "Milestone: Tanaka demo & UAT",
        ],
    )

    content_slide(
        prs,
        "Timeline",
        "Phase 3 — Polish & handoff (wks 8–10)",
        [
            "Hardening, performance pass, accessibility check",
            "Desired items where they fit (D1 / D2 / D3)",
            "CI integration with Meridian IT",
            "Final architecture handoff document",
            "Milestone: production cutover, warranty period begins",
        ],
    )

    content_slide(
        prs,
        "Investment",
        "Pricing",
        [
            "Fixed-fee:  $148,000",
            "Billed against four phase milestones, not hours",
            "Includes all required items (R1–R4) and Phase 3 desired items where they fit",
            "30-day warranty period post-cutover included",
            "Out of scope: net-new modules beyond Restocking, ERP integration, data migration",
        ],
    )

    content_slide(
        prs,
        "Why us",
        "What you get with this team",
        [
            "We leave codebases better-documented than we found them",
            "Tests-as-we-go, not tests-at-the-end",
            "Operator-centric design — your team stays in control",
            "Predictable weekly cadence, no week-eight surprises",
            "Direct line to senior engineers — no handoff to junior staff post-sale",
        ],
    )

    closing_slide(
        prs,
        "Let's make the next change a green build.",
        "Questions, clarifications, next steps — we're ready when you are.",
        "RFP MC-2026-0417  •  Response due May 8, 2026",
    )

    out = "capabilities-deck.pptx"
    prs.save(out)
    print(f"Wrote {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
